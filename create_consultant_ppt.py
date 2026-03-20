"""
AI 인보이스 처리 혁신 전략 - 컨설턴트급 PPT 생성 스크립트
맥킨지/BCG 스타일 디자인 시스템 적용
"""

import os
import tempfile
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
import numpy as np

# ─── 디자인 시스템 상수 ───────────────────────────────────────────────
COLOR = {
    'primary':   RGBColor(0x1A, 0x3A, 0x5F),  # Dark Navy
    'secondary': RGBColor(0x2E, 0x6C, 0xA4),  # Blue
    'accent':    RGBColor(0xD4, 0xA8, 0x43),  # Gold
    'success':   RGBColor(0x2E, 0x8B, 0x57),  # Green
    'danger':    RGBColor(0xC0, 0x39, 0x2B),  # Red
    'text':      RGBColor(0x2D, 0x2D, 0x2D),  # Dark Gray
    'light_bg':  RGBColor(0xF5, 0xF7, 0xFA),  # Light BG
    'white':     RGBColor(0xFF, 0xFF, 0xFF),
    'light_gray': RGBColor(0xE0, 0xE0, 0xE0),
    'medium_gray': RGBColor(0x99, 0x99, 0x99),
}

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.5)
TITLE_BAR_H = Inches(0.85)
FOOTER_H = Inches(0.35)
CONTENT_TOP = Inches(1.15)
CONTENT_BOTTOM = SLIDE_H - FOOTER_H - Inches(0.15)

FONT_KR = 'Malgun Gothic'
FONT_EN = 'Arial'

# matplotlib 한국어 폰트 설정
plt.rcParams['font.family'] = 'Malgun Gothic'
plt.rcParams['axes.unicode_minus'] = False

# 임시 차트 파일 관리
temp_files = []

def get_temp_path(name):
    path = os.path.join(tempfile.gettempdir(), f'ppt_chart_{name}.png')
    temp_files.append(path)
    return path


# ─── 헬퍼 함수 ───────────────────────────────────────────────────────

def set_font(run, size=11, bold=False, color=None, name=None, italic=False):
    """텍스트 run의 폰트를 설정"""
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    run.font.name = name or FONT_KR
    # 동아시아 폰트 설정
    from pptx.oxml.ns import qn
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', FONT_KR)


def add_title_bar(slide, title_text, subtitle_text=None):
    """네이비 타이틀 바 추가"""
    # 타이틀 바 배경
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, TITLE_BAR_H
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR['primary']
    bar.line.fill.background()

    # 타이틀 텍스트
    txBox = slide.shapes.add_textbox(MARGIN, Inches(0.12), SLIDE_W - 2 * MARGIN, Inches(0.55))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    set_font(run, size=24, bold=True, color=COLOR['white'])

    # 골드 악센트 라인
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN, TITLE_BAR_H - Inches(0.04), Inches(2), Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR['accent']
    line.line.fill.background()

    if subtitle_text:
        sub = slide.shapes.add_textbox(MARGIN, Inches(0.58), SLIDE_W - 2*MARGIN, Inches(0.25))
        stf = sub.text_frame
        sp = stf.paragraphs[0]
        sr = sp.add_run()
        sr.text = subtitle_text
        set_font(sr, size=11, color=COLOR['light_gray'])


def add_footer(slide, page_num):
    """하단 풋터 바"""
    footer_top = SLIDE_H - FOOTER_H
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, footer_top, SLIDE_W, FOOTER_H
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR['primary']
    bar.line.fill.background()

    # 회사명
    txBox = slide.shapes.add_textbox(MARGIN, footer_top + Inches(0.05), Inches(4), Inches(0.25))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "AI Invoice Automation Project  |  Confidential"
    set_font(run, size=9, color=COLOR['light_gray'])

    # 페이지 번호
    txBox2 = slide.shapes.add_textbox(SLIDE_W - Inches(1.5), footer_top + Inches(0.05), Inches(1), Inches(0.25))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    run2 = p2.add_run()
    run2.text = str(page_num)
    set_font(run2, size=9, color=COLOR['accent'], bold=True)


def add_card(slide, left, top, width, height, title, body_lines=None,
             icon_text=None, big_number=None, accent_color=None, fill_color=None):
    """KPI/정보 카드 생성"""
    ac = accent_color or COLOR['secondary']
    fc = fill_color or COLOR['white']

    # 카드 배경
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = fc
    card.line.color.rgb = COLOR['light_gray']
    card.line.width = Pt(0.5)
    # 둥근 모서리
    card.adjustments[0] = 0.05

    # 상단 악센트 라인
    accent_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.05)
    )
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = ac
    accent_line.line.fill.background()

    y_cursor = top + Inches(0.15)

    # 아이콘 텍스트 (큰 이모지/기호)
    if icon_text:
        ib = slide.shapes.add_textbox(left + Inches(0.15), y_cursor, width - Inches(0.3), Inches(0.4))
        itf = ib.text_frame
        ip = itf.paragraphs[0]
        ip.alignment = PP_ALIGN.CENTER
        ir = ip.add_run()
        ir.text = icon_text
        set_font(ir, size=22, color=ac)
        y_cursor += Inches(0.4)

    # 큰 숫자
    if big_number:
        nb = slide.shapes.add_textbox(left + Inches(0.15), y_cursor, width - Inches(0.3), Inches(0.45))
        ntf = nb.text_frame
        np_ = ntf.paragraphs[0]
        np_.alignment = PP_ALIGN.CENTER
        nr = np_.add_run()
        nr.text = big_number
        set_font(nr, size=28, bold=True, color=ac)
        y_cursor += Inches(0.45)

    # 제목
    tb = slide.shapes.add_textbox(left + Inches(0.15), y_cursor, width - Inches(0.3), Inches(0.35))
    ttf = tb.text_frame
    ttf.word_wrap = True
    tp = ttf.paragraphs[0]
    tp.alignment = PP_ALIGN.CENTER
    tr = tp.add_run()
    tr.text = title
    set_font(tr, size=12, bold=True, color=COLOR['text'])
    y_cursor += Inches(0.35)

    # 본문 라인
    if body_lines:
        bb = slide.shapes.add_textbox(left + Inches(0.15), y_cursor, width - Inches(0.3), height - (y_cursor - top) - Inches(0.15))
        btf = bb.text_frame
        btf.word_wrap = True
        for i, line in enumerate(body_lines):
            bp = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
            bp.alignment = PP_ALIGN.LEFT
            bp.space_before = Pt(2)
            br = bp.add_run()
            br.text = line
            set_font(br, size=10, color=COLOR['text'])


def add_process_arrow(slide, left, top, width, height, text, color=None):
    """프로세스 화살표 도형"""
    c = color or COLOR['secondary']
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.CHEVRON, left, top, width, height
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = c
    arrow.line.fill.background()
    tf = arrow.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    set_font(run, size=10, bold=True, color=COLOR['white'])
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    return arrow


def add_comparison_table(slide, left, top, width, headers, rows, col_widths=None):
    """비교 테이블 생성"""
    n_cols = len(headers)
    n_rows = len(rows) + 1  # +1 for header
    row_h = Inches(0.4)
    table_h = row_h * n_rows

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, table_h)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    # 헤더
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR['primary']
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                set_font(run, size=11, bold=True, color=COLOR['white'])

    # 데이터
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR['white'] if r % 2 == 0 else COLOR['light_bg']
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    set_font(run, size=10, color=COLOR['text'])

    return table_shape


# ─── matplotlib 차트 생성 ────────────────────────────────────────────

def create_bar_chart():
    """슬라이드 4: OCR 정확도 비교 막대 차트"""
    fig, ax = plt.subplots(figsize=(8, 4.5))
    fig.patch.set_facecolor('#F5F7FA')
    ax.set_facecolor('#F5F7FA')

    categories = ['문자 인식\n정확도', '테이블 구조\n인식률', '다국어\n지원', '처리 속도\n(장/분)']
    tesseract = [78, 45, 30, 15]
    claude = [99.5, 96, 95, 45]

    x = np.arange(len(categories))
    w = 0.35

    bars1 = ax.bar(x - w/2, tesseract, w, label='기존 OCR (Tesseract)',
                   color='#C0392B', alpha=0.85, edgecolor='white', linewidth=0.5)
    bars2 = ax.bar(x + w/2, claude, w, label='Claude Vision AI',
                   color='#2E6CA4', alpha=0.85, edgecolor='white', linewidth=0.5)

    # 값 라벨
    for bar in bars1:
        ax.text(bar.get_x() + bar.get_width()/2., bar.get_height() + 1.5,
                f'{bar.get_height():.0f}%' if bar.get_height() > 50 else f'{bar.get_height():.0f}',
                ha='center', va='bottom', fontsize=11, fontweight='bold', color='#C0392B')
    for bar in bars2:
        ax.text(bar.get_x() + bar.get_width()/2., bar.get_height() + 1.5,
                f'{bar.get_height():.1f}%' if bar.get_height() > 50 else f'{bar.get_height():.0f}',
                ha='center', va='bottom', fontsize=11, fontweight='bold', color='#2E6CA4')

    ax.set_xticks(x)
    ax.set_xticklabels(categories, fontsize=11)
    ax.set_ylim(0, 115)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color('#CCCCCC')
    ax.spines['bottom'].set_color('#CCCCCC')
    ax.tick_params(colors='#666666')
    ax.legend(loc='upper left', fontsize=11, framealpha=0.9)
    ax.set_ylabel('성능 (%/장)', fontsize=11, color='#666666')

    plt.tight_layout()
    path = get_temp_path('bar_chart')
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor=fig.get_facecolor())
    plt.close(fig)
    return path


def create_donut_chart():
    """슬라이드 8: 비용 절감 도넛 차트"""
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(8, 4))
    fig.patch.set_facecolor('#F5F7FA')

    # Before
    labels1 = ['인건비', 'IT유지보수', '오류비용', '기타']
    sizes1 = [2.5, 0.5, 0.8, 0.4]
    colors1 = ['#C0392B', '#E74C3C', '#F39C12', '#95A5A6']
    wedges1, texts1, autotexts1 = ax1.pie(
        sizes1, labels=labels1, colors=colors1, autopct='%1.0f%%',
        startangle=90, pctdistance=0.75, wedgeprops=dict(width=0.4, edgecolor='white'))
    for t in autotexts1:
        t.set_fontsize(9)
        t.set_fontweight('bold')
    for t in texts1:
        t.set_fontsize(9)
    ax1.set_title('AS-IS: 연간 ₩4.2억', fontsize=13, fontweight='bold', color='#C0392B', pad=15)

    # After
    labels2 = ['AI 운영비', 'IT유지보수', '인건비', '기타']
    sizes2 = [0.15, 0.1, 0.15, 0.1]
    colors2 = ['#2E6CA4', '#3498DB', '#2E8B57', '#95A5A6']
    wedges2, texts2, autotexts2 = ax2.pie(
        sizes2, labels=labels2, colors=colors2, autopct='%1.0f%%',
        startangle=90, pctdistance=0.75, wedgeprops=dict(width=0.4, edgecolor='white'))
    for t in autotexts2:
        t.set_fontsize(9)
        t.set_fontweight('bold')
    for t in texts2:
        t.set_fontsize(9)
    ax2.set_title('TO-BE: 연간 ₩0.5억', fontsize=13, fontweight='bold', color='#2E6CA4', pad=15)

    plt.tight_layout()
    path = get_temp_path('donut_chart')
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor=fig.get_facecolor())
    plt.close(fig)
    return path


def create_timeline_chart():
    """슬라이드 14: 6개월 4단계 타임라인"""
    fig, ax = plt.subplots(figsize=(10, 3.5))
    fig.patch.set_facecolor('#F5F7FA')
    ax.set_facecolor('#F5F7FA')

    phases = [
        ('Phase 1\n기반 구축', '1~2개월', '#1A3A5F'),
        ('Phase 2\nAI 통합', '2~3개월', '#2E6CA4'),
        ('Phase 3\nERP 연동', '4~5개월', '#D4A843'),
        ('Phase 4\n전사 확대', '5~6개월', '#2E8B57'),
    ]

    y = 0.5
    bar_h = 0.6
    starts = [0, 1.5, 3.5, 4.5]
    durations = [2, 2, 2, 2]

    for i, ((label, period, color), start, dur) in enumerate(zip(phases, starts, durations)):
        ax.barh(y, dur, left=start, height=bar_h, color=color, alpha=0.9,
                edgecolor='white', linewidth=2)
        ax.text(start + dur/2, y, f'{label}\n({period})', ha='center', va='center',
                fontsize=10, fontweight='bold', color='white')

    ax.set_xlim(-0.5, 7)
    ax.set_ylim(-0.2, 1.2)
    ax.set_xticks(range(7))
    ax.set_xticklabels(['M1', 'M2', 'M3', 'M4', 'M5', 'M6', ''], fontsize=10, color='#666666')
    ax.set_yticks([])
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_visible(False)
    ax.spines['bottom'].set_color('#CCCCCC')

    # 마일스톤 마커
    milestones = [(2, 'MVP\n완료'), (3.5, 'AI 파이프라인\n구축 완료'), (5, 'ERP 연동\n완료'), (6.5, 'Go-Live')]
    for mx, mlabel in milestones:
        ax.plot(mx, 1.0, 'v', markersize=10, color='#D4A843')
        ax.text(mx, 1.1, mlabel, ha='center', va='bottom', fontsize=8, color='#1A3A5F', fontweight='bold')

    plt.tight_layout()
    path = get_temp_path('timeline')
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor=fig.get_facecolor())
    plt.close(fig)
    return path


# ─── 슬라이드 생성 함수 ──────────────────────────────────────────────

def slide_01_cover(prs):
    """표지: 인보이스 처리 혁신 전략"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # 그라데이션 배경 - 네이비 풀 커버
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR['primary']
    bg.line.fill.background()

    # 상단 장식 라인
    deco = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.08)
    )
    deco.fill.solid()
    deco.fill.fore_color.rgb = COLOR['accent']
    deco.line.fill.background()

    # 골드 악센트 라인 (중앙)
    gold_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2.7), Inches(2), Inches(0.05)
    )
    gold_line.fill.solid()
    gold_line.fill.fore_color.rgb = COLOR['accent']
    gold_line.line.fill.background()

    # 메인 타이틀
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.9), Inches(10), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "인보이스 처리 혁신 전략"
    set_font(run, size=40, bold=True, color=COLOR['white'])

    # 서브 타이틀
    p2 = tf.add_paragraph()
    p2.space_before = Pt(12)
    run2 = p2.add_run()
    run2.text = "AI 기반 자동화로 연간 ₩3.7억 절감 및 처리 시간 97% 단축"
    set_font(run2, size=18, color=COLOR['accent'])

    # 하단 정보
    info_box = slide.shapes.add_textbox(Inches(1.5), Inches(5.5), Inches(6), Inches(1))
    itf = info_box.text_frame
    itf.word_wrap = True
    lines = ["AI Invoice Automation Project", "2026년 3월  |  전략 기획팀"]
    for i, line in enumerate(lines):
        ip = itf.paragraphs[0] if i == 0 else itf.add_paragraph()
        ip.space_before = Pt(4)
        ir = ip.add_run()
        ir.text = line
        set_font(ir, size=13, color=COLOR['light_gray'])

    # 하단 골드 라인
    bottom_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08)
    )
    bottom_line.fill.solid()
    bottom_line.fill.fore_color.rgb = COLOR['accent']
    bottom_line.line.fill.background()

    # 우측 장식 - 큰 반투명 원
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(9), Inches(1.5), Inches(4.5), Inches(4.5)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0x25, 0x50, 0x80)
    circle.line.fill.background()

    circle2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(9.8), Inches(2.3), Inches(3), Inches(3))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = RGBColor(0x30, 0x60, 0x90)
    circle2.line.fill.background()


def slide_02_problem(prs):
    """수작업 인보이스 처리로 연간 ₩3.7억이 낭비되고 있습니다"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "수작업 인보이스 처리로 연간 ₩3.7억이 낭비되고 있습니다")
    add_footer(slide, 2)

    # 4 KPI 카드
    cards = [
        ("처리 시간", "45분/건", "수작업 데이터 입력\n검증에 평균 45분 소요", COLOR['danger'], ">>"),
        ("오류율", "12.5%", "수동 입력 시\n평균 오류 발생률", COLOR['danger'], "!!"),
        ("월 처리량", "2,400건", "3명의 담당자가\n월 처리하는 인보이스 수", COLOR['secondary'], "##"),
        ("연간 비용", "₩4.2억", "인건비+IT유지보수\n+오류 처리 비용 합계", COLOR['danger'], "$$"),
    ]

    card_w = Inches(2.8)
    card_h = Inches(2.8)
    start_x = Inches(0.65)
    gap = Inches(0.25)
    y = CONTENT_TOP + Inches(0.2)

    for i, (title, number, body, color, icon) in enumerate(cards):
        x = start_x + i * (card_w + gap)
        add_card(slide, x, y, card_w, card_h, title,
                 body_lines=body.split('\n'), icon_text=icon,
                 big_number=number, accent_color=color)

    # 하단 손실 합계 바
    bar_y = y + card_h + Inches(0.4)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), bar_y, Inches(12), Inches(0.8)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0xFD, 0xED, 0xED)
    bar.line.color.rgb = COLOR['danger']
    bar.line.width = Pt(1.5)
    bar.adjustments[0] = 0.15

    txBox = slide.shapes.add_textbox(Inches(1), bar_y + Inches(0.15), Inches(11), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "연간 총 손실: ₩3.7억  =  인건비 ₩2.5억  +  오류 처리 ₩0.8억  +  기회 비용 ₩0.4억"
    set_font(run, size=14, bold=True, color=COLOR['danger'])


def slide_03_solution(prs):
    """AI 기반 5단계 자동화로 처리 시간을 97% 단축합니다"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "AI 기반 5단계 자동화로 처리 시간을 97% 단축합니다")
    add_footer(slide, 3)

    # 5단계 프로세스 플로우
    steps = [
        ("1. 수신", "이메일/스캔\n자동 수집"),
        ("2. 인식", "Claude Vision\nAI 추출"),
        ("3. 검증", "3단계\n자동 검증"),
        ("4. 승인", "규칙 기반\n자동 결재"),
        ("5. 전기", "ERP 자동\n전표 생성"),
    ]

    arrow_w = Inches(2.2)
    arrow_h = Inches(1.0)
    start_x = Inches(0.5)
    gap = Inches(0.25)
    y = CONTENT_TOP + Inches(0.6)

    colors = [COLOR['secondary'], COLOR['primary'], COLOR['accent'],
              COLOR['success'], COLOR['primary']]

    for i, ((title, desc), color) in enumerate(zip(steps, colors)):
        x = start_x + i * (arrow_w + gap)
        add_process_arrow(slide, x, y, arrow_w, arrow_h, title, color)

        # 설명 텍스트 (아래)
        tb = slide.shapes.add_textbox(x, y + arrow_h + Inches(0.1), arrow_w, Inches(0.6))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = desc
        set_font(run, size=10, color=COLOR['text'])

    # 하단 결과 카드
    result_y = y + arrow_h + Inches(1.2)

    results = [
        ("처리 시간", "45분 → 1.5분", "(97% 단축)", COLOR['success']),
        ("오류율", "12.5% → 0.3%", "(97.6% 감소)", COLOR['success']),
        ("자동화율", "0% → 94%", "(완전 자동화)", COLOR['secondary']),
    ]

    rw = Inches(3.5)
    rh = Inches(1.2)
    rx = Inches(0.9)
    for i, (title, value, sub, color) in enumerate(results):
        x = rx + i * (rw + Inches(0.4))
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, result_y, rw, rh
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR['white']
        card.line.color.rgb = color
        card.line.width = Pt(1.5)
        card.adjustments[0] = 0.1

        tb = slide.shapes.add_textbox(x + Inches(0.2), result_y + Inches(0.12), rw - Inches(0.4), Inches(0.25))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = title
        set_font(r, size=11, color=COLOR['medium_gray'])

        tb2 = slide.shapes.add_textbox(x + Inches(0.2), result_y + Inches(0.4), rw - Inches(0.4), Inches(0.35))
        tf2 = tb2.text_frame
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = value
        set_font(r2, size=18, bold=True, color=color)

        tb3 = slide.shapes.add_textbox(x + Inches(0.2), result_y + Inches(0.78), rw - Inches(0.4), Inches(0.25))
        tf3 = tb3.text_frame
        p3 = tf3.paragraphs[0]
        p3.alignment = PP_ALIGN.CENTER
        r3 = p3.add_run()
        r3.text = sub
        set_font(r3, size=10, color=color)


def slide_04_ocr_comparison(prs):
    """Claude Vision은 기존 OCR 대비 오류를 95% 감소시킵니다"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Claude Vision은 기존 OCR 대비 오류를 95% 감소시킵니다")
    add_footer(slide, 4)

    # matplotlib 차트 삽입
    chart_path = create_bar_chart()
    slide.shapes.add_picture(chart_path, Inches(0.6), CONTENT_TOP + Inches(0.2),
                             Inches(7.5), Inches(4.2))

    # 우측 핵심 포인트
    points = [
        ("99.5%", "문자 인식 정확도", "한국어/영문/일문 혼합 문서도\n정확하게 인식"),
        ("96%", "테이블 구조 인식", "복잡한 인보이스 테이블 구조를\n자동으로 파싱"),
        ("3배", "처리 속도 향상", "기존 OCR 대비 3배 빠른\n실시간 처리"),
    ]

    px = Inches(8.5)
    pw = Inches(4.2)
    ph = Inches(1.3)
    py = CONTENT_TOP + Inches(0.3)

    for i, (num, title, desc) in enumerate(points):
        y = py + i * (ph + Inches(0.15))
        add_card(slide, px, y, pw, ph, title,
                 body_lines=desc.split('\n'), big_number=num,
                 accent_color=COLOR['secondary'])


def slide_05_validation(prs):
    """3단계 검증으로 오류와 사기를 원천 차단합니다"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "3단계 검증으로 오류와 사기를 원천 차단합니다")
    add_footer(slide, 5)

    # 3컬럼 카드
    validations = [
        ("1단계: 데이터 검증", COLOR['secondary'],
         ["필수 필드 완전성 확인", "데이터 형식/범위 검증", "사업자번호 유효성 체크",
          "날짜/금액 형식 표준화"]),
        ("2단계: 비즈니스 룰", COLOR['accent'],
         ["거래처 마스터 데이터 매칭", "계약 단가 vs 청구 단가 비교", "수량 x 단가 = 합계 교차검증",
          "중복 인보이스 탐지"]),
        ("3단계: AI 이상 탐지", COLOR['danger'],
         ["과거 거래 패턴 이탈 감지", "비정상 금액 자동 플래그", "사기 패턴 머신러닝 탐지",
          "실시간 리스크 스코어링"]),
    ]

    card_w = Inches(3.8)
    card_h = Inches(4.2)
    start_x = Inches(0.55)
    gap = Inches(0.35)
    y = CONTENT_TOP + Inches(0.3)

    for i, (title, color, items) in enumerate(validations):
        x = start_x + i * (card_w + gap)

        # 카드 배경
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR['white']
        card.line.color.rgb = COLOR['light_gray']
        card.line.width = Pt(0.5)
        card.adjustments[0] = 0.05

        # 상단 색상 바
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y, card_w, Inches(0.06)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = color
        accent.line.fill.background()

        # 단계 번호 원
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x + Inches(0.2), y + Inches(0.25), Inches(0.6), Inches(0.6)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        ctf = circle.text_frame
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = str(i + 1)
        set_font(cr, size=20, bold=True, color=COLOR['white'])
        ctf.paragraphs[0].space_before = Pt(5)

        # 제목
        tb = slide.shapes.add_textbox(x + Inches(0.95), y + Inches(0.33), card_w - Inches(1.2), Inches(0.4))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        set_font(r, size=14, bold=True, color=color)

        # 체크리스트 아이템
        items_box = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(1.1), card_w - Inches(0.6), Inches(3))
        itf = items_box.text_frame
        itf.word_wrap = True
        for j, item in enumerate(items):
            ip = itf.paragraphs[0] if j == 0 else itf.add_paragraph()
            ip.space_before = Pt(8)
            ir = ip.add_run()
            ir.text = f"  {item}"
            set_font(ir, size=11, color=COLOR['text'])


def slide_06_value_prop(prs):
    """4가지 핵심 가치가 투자 대비 확실한 수익을 보장합니다"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "4가지 핵심 가치가 투자 대비 확실한 수익을 보장합니다")
    add_footer(slide, 6)

    values = [
        ("$$", "비용 절감", "₩3.7억/년", "인건비 88% 절감\n오류 처리 비용 95% 감소\n운영 효율 극대화", COLOR['success']),
        (">>", "속도 혁신", "97% 단축", "45분 → 1.5분 처리\n실시간 ERP 반영\n24/7 자동 처리", COLOR['secondary']),
        ("!!", "정확도 향상", "99.5%", "AI 기반 데이터 추출\n3단계 자동 검증\n오류율 0.3% 이하", COLOR['primary']),
        ("##", "확장성", "무제한", "클라우드 기반 확장\n다국어 인보이스 지원\n신규 거래처 즉시 적용", COLOR['accent']),
    ]

    card_w = Inches(2.85)
    card_h = Inches(3.8)
    start_x = Inches(0.5)
    gap = Inches(0.2)
    y = CONTENT_TOP + Inches(0.3)

    for i, (icon, title, number, desc, color) in enumerate(values):
        x = start_x + i * (card_w + gap)
        add_card(slide, x, y, card_w, card_h, title,
                 body_lines=desc.split('\n'), icon_text=icon,
                 big_number=number, accent_color=color)


def slide_07_as_is_to_be(prs):
    """업무 프로세스가 근본적으로 바뀝니다"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "업무 프로세스가 근본적으로 바뀝니다")
    add_footer(slide, 7)

    # AS-IS (좌측)
    as_is_x = Inches(0.5)
    as_is_w = Inches(5.9)
    y = CONTENT_TOP + Inches(0.15)

    # AS-IS 헤더
    as_header = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, as_is_x, y, as_is_w, Inches(0.55)
    )
    as_header.fill.solid()
    as_header.fill.fore_color.rgb = COLOR['danger']
    as_header.line.fill.background()
    as_header.adjustments[0] = 0.15
    tb = slide.shapes.add_textbox(as_is_x + Inches(0.2), y + Inches(0.1), as_is_w - Inches(0.4), Inches(0.35))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "AS-IS  |  현재 프로세스"
    set_font(r, size=16, bold=True, color=COLOR['white'])

    as_items = [
        ("수동 수신", "이메일 확인 → 첨부파일 다운로드 → 폴더 분류"),
        ("수동 입력", "인보이스 내용을 ERP에 하나씩 수동 입력"),
        ("수동 검증", "담당자가 육안으로 데이터 확인 및 교차 검증"),
        ("수동 결재", "결재 라인별 이메일/메신저 요청 → 대기"),
        ("수동 전기", "승인 후 ERP에 수동 전표 생성 및 입력"),
    ]

    for i, (step, desc) in enumerate(as_items):
        iy = y + Inches(0.7) + i * Inches(0.78)
        # 번호 원
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, as_is_x + Inches(0.15), iy + Inches(0.08), Inches(0.4), Inches(0.4)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLOR['danger']
        circle.line.fill.background()
        ctf = circle.text_frame
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = str(i + 1)
        set_font(cr, size=12, bold=True, color=COLOR['white'])

        # 텍스트
        tb = slide.shapes.add_textbox(as_is_x + Inches(0.7), iy, as_is_w - Inches(0.9), Inches(0.6))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = step
        set_font(r, size=12, bold=True, color=COLOR['danger'])
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = desc
        set_font(r2, size=10, color=COLOR['text'])

    # 중앙 화살표
    arrow_x = Inches(6.55)
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, arrow_x, Inches(3.2), Inches(0.6), Inches(1.0)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = COLOR['accent']
    arrow.line.fill.background()

    # TO-BE (우측)
    to_be_x = Inches(7.3)
    to_be_w = Inches(5.5)

    to_header = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, to_be_x, y, to_be_w, Inches(0.55)
    )
    to_header.fill.solid()
    to_header.fill.fore_color.rgb = COLOR['success']
    to_header.line.fill.background()
    to_header.adjustments[0] = 0.15
    tb = slide.shapes.add_textbox(to_be_x + Inches(0.2), y + Inches(0.1), to_be_w - Inches(0.4), Inches(0.35))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "TO-BE  |  AI 자동화 프로세스"
    set_font(r, size=16, bold=True, color=COLOR['white'])

    to_items = [
        ("자동 수신", "이메일 모니터링 → 자동 분류 → 큐 등록"),
        ("AI 인식", "Claude Vision이 자동으로 데이터 추출"),
        ("자동 검증", "3단계 자동 검증 + AI 이상 탐지"),
        ("자동 승인", "금액별 자동 승인 규칙 적용"),
        ("자동 전기", "ERP API 연동 자동 전표 생성"),
    ]

    for i, (step, desc) in enumerate(to_items):
        iy = y + Inches(0.7) + i * Inches(0.78)
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, to_be_x + Inches(0.15), iy + Inches(0.08), Inches(0.4), Inches(0.4)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLOR['success']
        circle.line.fill.background()
        ctf = circle.text_frame
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = str(i + 1)
        set_font(cr, size=12, bold=True, color=COLOR['white'])

        tb = slide.shapes.add_textbox(to_be_x + Inches(0.7), iy, to_be_w - Inches(0.9), Inches(0.6))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = step
        set_font(r, size=12, bold=True, color=COLOR['success'])
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = desc
        set_font(r2, size=10, color=COLOR['text'])


def slide_08_roi(prs):
    """도입 후 연간 ₩3.7억 절감, 18개월 내 ROI 달성"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "도입 후 연간 ₩3.7억 절감, 18개월 내 ROI 달성")
    add_footer(slide, 8)

    # 4개 메트릭 카드 (상단)
    metrics = [
        ("연간 절감액", "₩3.7억", COLOR['success']),
        ("ROI 달성", "18개월", COLOR['secondary']),
        ("투자 회수율", "340%", COLOR['accent']),
        ("운영비 절감", "88%", COLOR['primary']),
    ]

    card_w = Inches(2.8)
    card_h = Inches(1.5)
    start_x = Inches(0.65)
    gap = Inches(0.25)
    y = CONTENT_TOP + Inches(0.15)

    for i, (title, number, color) in enumerate(metrics):
        x = start_x + i * (card_w + gap)
        add_card(slide, x, y, card_w, card_h, title,
                 big_number=number, accent_color=color)

    # 도넛 차트
    chart_path = create_donut_chart()
    slide.shapes.add_picture(chart_path, Inches(0.65), y + card_h + Inches(0.3),
                             Inches(11.5), Inches(3.3))


def slide_09_approval(prs):
    """금액별 자동 승인 규칙으로 결재 병목을 제거합니다"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "금액별 자동 승인 규칙으로 결재 병목을 제거합니다")
    add_footer(slide, 9)

    # 계단형 승인 플로우
    levels = [
        ("₩100만 미만", "자동 승인", "AI 검증 통과 시\n즉시 자동 승인", COLOR['success'], Inches(0.7)),
        ("₩100만~₩500만", "팀장 승인", "팀장 1인 결재\n(모바일 알림)", COLOR['secondary'], Inches(1.6)),
        ("₩500만~₩2,000만", "부서장 승인", "팀장 → 부서장\n순차 결재", COLOR['accent'], Inches(2.5)),
        ("₩2,000만 이상", "임원 승인", "팀장 → 부서장 → 임원\n3단계 결재", COLOR['primary'], Inches(3.4)),
    ]

    step_w = Inches(2.8)
    step_h = Inches(1.8)
    start_x = Inches(0.6)
    gap = Inches(0.2)

    for i, (amount, approver, desc, color, y_off) in enumerate(levels):
        x = start_x + i * (step_w + gap)
        y = CONTENT_TOP + y_off

        # 카드
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, step_w, step_h
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR['white']
        card.line.color.rgb = color
        card.line.width = Pt(2)
        card.adjustments[0] = 0.08

        # 상단 바
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, step_w, Inches(0.06))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        # 금액 범위
        tb = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.15), step_w - Inches(0.2), Inches(0.3))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = amount
        set_font(r, size=11, bold=True, color=color)

        # 승인자
        tb2 = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.45), step_w - Inches(0.2), Inches(0.35))
        tf2 = tb2.text_frame
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = approver
        set_font(r2, size=16, bold=True, color=COLOR['text'])

        # 설명
        tb3 = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.9), step_w - Inches(0.2), Inches(0.8))
        tf3 = tb3.text_frame
        tf3.word_wrap = True
        p3 = tf3.paragraphs[0]
        p3.alignment = PP_ALIGN.CENTER
        r3 = p3.add_run()
        r3.text = desc
        set_font(r3, size=10, color=COLOR['medium_gray'])

        # 연결 화살표 (마지막 제외)
        if i < len(levels) - 1:
            arr = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, x + step_w + Inches(0.02), y + Inches(0.6),
                Inches(0.16), Inches(0.4)
            )
            arr.fill.solid()
            arr.fill.fore_color.rgb = COLOR['light_gray']
            arr.line.fill.background()

    # 하단 설명
    note_y = CONTENT_TOP + Inches(5.2)
    note = slide.shapes.add_textbox(Inches(0.6), note_y, Inches(12), Inches(0.4))
    ntf = note.text_frame
    np_ = ntf.paragraphs[0]
    np_.alignment = PP_ALIGN.CENTER
    nr = np_.add_run()
    nr.text = "* 전체 인보이스의 67%가 ₩100만 미만으로 자동 승인 대상  →  결재 대기 시간 평균 2일 → 0일"
    set_font(nr, size=11, bold=True, color=COLOR['secondary'])


def slide_10_build_vs_buy(prs):
    """자체 구축이 SaaS 대비 TCO와 유연성에서 우위입니다"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "자체 구축이 SaaS 대비 TCO와 유연성에서 우위입니다")
    add_footer(slide, 10)

    # 비교 테이블
    headers = ["항목", "자체 구축 (권장)", "SaaS 솔루션"]
    rows = [
        ["초기 비용", "₩1.5억 (일시)", "₩0.3억/년 (구독)"],
        ["3년 TCO", "₩2.1억", "₩3.6억"],
        ["커스터마이징", "무제한", "제한적"],
        ["데이터 보안", "온프레미스 보관", "외부 서버"],
        ["ERP 연동", "직접 API 연동", "제한적 연동"],
        ["AI 모델 선택", "자유 (Claude 등)", "벤더 종속"],
        ["확장성", "자유로운 확장", "플랜 종속"],
    ]

    table_shape = add_comparison_table(
        slide, Inches(0.5), CONTENT_TOP + Inches(0.2), Inches(12.3),
        headers, rows,
        col_widths=[Inches(2.5), Inches(4.9), Inches(4.9)]
    )

    # 자체 구축 열 강조 (2번째 열 각 셀에 체크 표시 효과)
    table = table_shape.table
    for r in range(1, len(rows) + 1):
        cell = table.cell(r, 1)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)

    # 하단 4개 장점 카드
    advantages = [
        ("데이터 주권", "민감 데이터\n사내 보관"),
        ("비용 효율", "3년 TCO\n42% 절감"),
        ("맞춤 개발", "업무 특화\n기능 구현"),
        ("기술 내재화", "AI 역량\n조직 축적"),
    ]

    card_w = Inches(2.8)
    card_h = Inches(1.3)
    start_x = Inches(0.65)
    gap = Inches(0.25)
    card_y = CONTENT_TOP + Inches(3.8)

    for i, (title, desc) in enumerate(advantages):
        x = start_x + i * (card_w + gap)
        add_card(slide, x, card_y, card_w, card_h, title,
                 body_lines=desc.split('\n'), accent_color=COLOR['success'])


def slide_11_ai_erp(prs):
    """AI Agent가 처리하고 ERP가 자동 실행합니다"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "AI Agent가 처리하고 ERP가 자동 실행합니다")
    add_footer(slide, 11)

    # 상단 프로세스 플로우
    flow_steps = [
        ("인보이스\n수신", COLOR['secondary']),
        ("Claude\nVision AI", COLOR['primary']),
        ("데이터\n검증/변환", COLOR['accent']),
        ("결재\n워크플로우", COLOR['success']),
        ("ERP\n자동 전기", COLOR['primary']),
    ]

    arrow_w = Inches(2.2)
    arrow_h = Inches(0.85)
    start_x = Inches(0.5)
    gap = Inches(0.2)
    y = CONTENT_TOP + Inches(0.2)

    for i, (text, color) in enumerate(flow_steps):
        x = start_x + i * (arrow_w + gap)
        add_process_arrow(slide, x, y, arrow_w, arrow_h, text, color)

    # 하단 ERP 연동 카드
    erp_systems = [
        ("SAP S/4HANA", "RFC/BAPI 연동\nFI 모듈 자동 전기\n실시간 데이터 동기화"),
        ("Oracle EBS", "REST API 연동\nAP 모듈 자동 처리\n멀티 통화 지원"),
        ("더존 iCUBE", "WebService 연동\n회계 전표 자동 생성\n국내 세무 최적화"),
        ("MS Dynamics", "OData API 연동\n매입 전표 자동화\n글로벌 법인 지원"),
    ]

    card_w = Inches(2.85)
    card_h = Inches(2.8)
    start_x = Inches(0.5)
    gap = Inches(0.2)
    card_y = CONTENT_TOP + Inches(1.5)

    colors = [COLOR['secondary'], COLOR['danger'], COLOR['success'], COLOR['accent']]
    for i, ((title, desc), color) in enumerate(zip(erp_systems, colors)):
        x = start_x + i * (card_w + gap)
        add_card(slide, x, card_y, card_w, card_h, title,
                 body_lines=desc.split('\n'), accent_color=color)

    # 하단 통합 메시지
    msg_y = card_y + card_h + Inches(0.3)
    msg_bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), msg_y, Inches(12.3), Inches(0.65)
    )
    msg_bar.fill.solid()
    msg_bar.fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xFE)
    msg_bar.line.color.rgb = COLOR['secondary']
    msg_bar.line.width = Pt(1)
    msg_bar.adjustments[0] = 0.15

    tb = slide.shapes.add_textbox(Inches(0.8), msg_y + Inches(0.12), Inches(11.5), Inches(0.4))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "AI Agent가 인보이스를 자율적으로 처리하고, 검증 완료 후 ERP에 자동 전기합니다"
    set_font(r, size=13, bold=True, color=COLOR['primary'])


def slide_12_architecture(prs):
    """검증된 기술 스택으로 확장 가능한 아키텍처를 구축합니다"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "검증된 기술 스택으로 확장 가능한 아키텍처를 구축합니다")
    add_footer(slide, 12)

    # 4계층 아키텍처
    layers = [
        ("프레젠테이션 계층", "React 18  |  TypeScript  |  Tailwind CSS  |  Chart.js",
         COLOR['secondary'], Inches(0.2)),
        ("API / 비즈니스 로직 계층", "FastAPI (Python)  |  JWT Auth  |  Rate Limiting  |  WebSocket",
         COLOR['primary'], Inches(1.5)),
        ("AI / ML 계층", "Claude Vision API  |  LangChain  |  Vector DB  |  Custom Models",
         COLOR['accent'], Inches(2.8)),
        ("데이터 / 인프라 계층", "PostgreSQL  |  Redis Cache  |  MinIO (S3)  |  Docker / K8s",
         COLOR['success'], Inches(4.1)),
    ]

    layer_w = Inches(12.3)
    layer_h = Inches(1.1)
    start_x = Inches(0.5)

    for title, tech, color, y_off in layers:
        y = CONTENT_TOP + y_off

        # 계층 배경
        layer = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, start_x, y, layer_w, layer_h
        )
        layer.fill.solid()
        layer.fill.fore_color.rgb = COLOR['white']
        layer.line.color.rgb = color
        layer.line.width = Pt(2)
        layer.adjustments[0] = 0.08

        # 좌측 색상 바
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, start_x, y, Inches(0.08), layer_h
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        # 계층 이름
        tb = slide.shapes.add_textbox(start_x + Inches(0.3), y + Inches(0.1), Inches(5), Inches(0.35))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        set_font(r, size=14, bold=True, color=color)

        # 기술 스택
        tb2 = slide.shapes.add_textbox(start_x + Inches(0.3), y + Inches(0.5), layer_w - Inches(0.6), Inches(0.4))
        tf2 = tb2.text_frame
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = tech
        set_font(r2, size=12, color=COLOR['text'])

        # 계층 간 화살표
        if y_off < Inches(4.0):
            arr = slide.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW, Inches(6.5), y + layer_h, Inches(0.35), Inches(0.35)
            )
            arr.fill.solid()
            arr.fill.fore_color.rgb = COLOR['light_gray']
            arr.line.fill.background()


def slide_13_security(prs):
    """엔터프라이즈급 보안으로 컴플라이언스를 충족합니다"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "엔터프라이즈급 보안으로 컴플라이언스를 충족합니다")
    add_footer(slide, 13)

    security_items = [
        ("데이터 암호화", COLOR['primary'],
         ["AES-256 저장 데이터 암호화", "TLS 1.3 전송 구간 암호화",
          "키 관리 시스템 (KMS)", "인보이스 이미지 자동 마스킹"]),
        ("접근 제어", COLOR['secondary'],
         ["RBAC 역할 기반 접근 제어", "JWT + OAuth 2.0 인증",
          "IP 화이트리스트", "세션 타임아웃 관리"]),
        ("감사 추적", COLOR['accent'],
         ["모든 조회/수정 이력 기록", "변경 불가 감사 로그",
          "실시간 이상 행위 감지", "월간 보안 리포트 자동 생성"]),
        ("컴플라이언스", COLOR['success'],
         ["전자세금계산서 법적 요건", "개인정보보호법 (PIPA) 준수",
          "내부회계관리제도 대응", "SOX / K-IFRS 감사 지원"]),
    ]

    card_w = Inches(2.85)
    card_h = Inches(4.2)
    start_x = Inches(0.5)
    gap = Inches(0.2)
    y = CONTENT_TOP + Inches(0.3)

    for i, (title, color, items) in enumerate(security_items):
        x = start_x + i * (card_w + gap)

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR['white']
        card.line.color.rgb = COLOR['light_gray']
        card.line.width = Pt(0.5)
        card.adjustments[0] = 0.05

        # 상단 바
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, Inches(0.06))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        # 아이콘 원
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x + card_w/2 - Inches(0.3), y + Inches(0.2), Inches(0.6), Inches(0.6)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        ctf = circle.text_frame
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = str(i + 1)
        set_font(cr, size=18, bold=True, color=COLOR['white'])

        # 제목
        tb = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.95), card_w - Inches(0.2), Inches(0.35))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = title
        set_font(r, size=14, bold=True, color=color)

        # 아이템
        items_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(1.45), card_w - Inches(0.4), Inches(2.6))
        itf = items_box.text_frame
        itf.word_wrap = True
        for j, item in enumerate(items):
            ip = itf.paragraphs[0] if j == 0 else itf.add_paragraph()
            ip.space_before = Pt(8)
            ir = ip.add_run()
            ir.text = f"  {item}"
            set_font(ir, size=10, color=COLOR['text'])


def slide_14_roadmap(prs):
    """6개월 4단계 로드맵으로 단계적 전환합니다"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "6개월 4단계 로드맵으로 단계적 전환합니다")
    add_footer(slide, 14)

    # 타임라인 차트
    chart_path = create_timeline_chart()
    slide.shapes.add_picture(chart_path, Inches(0.5), CONTENT_TOP + Inches(0.1),
                             Inches(12.3), Inches(2.8))

    # 하단 4단계 상세 카드
    phases = [
        ("Phase 1", "기반 구축", "1~2개월",
         ["개발 환경 구축", "DB 설계/마이그레이션", "인증/권한 시스템", "기본 UI 프레임워크"], COLOR['primary']),
        ("Phase 2", "AI 통합", "2~3개월",
         ["Claude Vision 연동", "데이터 추출 파이프라인", "3단계 검증 엔진", "대시보드 구현"], COLOR['secondary']),
        ("Phase 3", "ERP 연동", "4~5개월",
         ["ERP API 개발", "승인 워크플로우", "자동 전기 시스템", "통합 테스트"], COLOR['accent']),
        ("Phase 4", "전사 확대", "5~6개월",
         ["파일럿 운영", "사용자 교육", "전사 롤아웃", "모니터링/최적화"], COLOR['success']),
    ]

    card_w = Inches(2.85)
    card_h = Inches(2.8)
    start_x = Inches(0.5)
    gap = Inches(0.2)
    card_y = CONTENT_TOP + Inches(3.0)

    for i, (phase, title, period, items, color) in enumerate(phases):
        x = start_x + i * (card_w + gap)

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, card_y, card_w, card_h
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR['white']
        card.line.color.rgb = color
        card.line.width = Pt(1.5)
        card.adjustments[0] = 0.06

        # 상단 바
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, card_y, card_w, Inches(0.06))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        # Phase 번호
        tb = slide.shapes.add_textbox(x + Inches(0.1), card_y + Inches(0.15), card_w - Inches(0.2), Inches(0.25))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = f"{phase}  |  {period}"
        set_font(r, size=10, bold=True, color=color)

        # 제목
        tb2 = slide.shapes.add_textbox(x + Inches(0.1), card_y + Inches(0.42), card_w - Inches(0.2), Inches(0.35))
        tf2 = tb2.text_frame
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = title
        set_font(r2, size=15, bold=True, color=COLOR['text'])

        # 아이템
        items_box = slide.shapes.add_textbox(x + Inches(0.2), card_y + Inches(0.9), card_w - Inches(0.4), Inches(1.8))
        itf = items_box.text_frame
        itf.word_wrap = True
        for j, item in enumerate(items):
            ip = itf.paragraphs[0] if j == 0 else itf.add_paragraph()
            ip.space_before = Pt(6)
            ir = ip.add_run()
            ir.text = f"  {item}"
            set_font(ir, size=10, color=COLOR['text'])


def slide_15_differentiators(prs):
    """경쟁사가 따라올 수 없는 4가지 차별화 요소"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "경쟁사가 따라올 수 없는 4가지 차별화 요소")
    add_footer(slide, 15)

    diffs = [
        ("01", "Claude Vision AI", COLOR['primary'],
         ["최신 멀티모달 AI 엔진 적용", "한국어 특화 인식 정확도 99.5%",
          "비정형 문서도 구조화 가능", "지속적 모델 업데이트"]),
        ("02", "자율형 AI Agent", COLOR['secondary'],
         ["규칙 기반 + AI 판단 하이브리드", "상황별 자동 에스컬레이션",
          "학습 기반 예외 처리 개선", "멀티 에이전트 협업"]),
        ("03", "실시간 ERP 연동", COLOR['accent'],
         ["4대 ERP 네이티브 연동", "양방향 실시간 동기화",
          "트랜잭션 무결성 보장", "제로 다운타임 운영"]),
        ("04", "엔드투엔드 자동화", COLOR['success'],
         ["수신→인식→검증→승인→전기", "인적 개입 최소화 (6%)",
          "예외 건만 사람이 처리", "전 과정 감사 추적"]),
    ]

    card_w = Inches(2.85)
    card_h = Inches(4.5)
    start_x = Inches(0.5)
    gap = Inches(0.2)
    y = CONTENT_TOP + Inches(0.2)

    for i, (num, title, color, items) in enumerate(diffs):
        x = start_x + i * (card_w + gap)

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR['white']
        card.line.color.rgb = COLOR['light_gray']
        card.line.width = Pt(0.5)
        card.adjustments[0] = 0.05

        # 상단 바
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, Inches(0.06))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        # 번호
        tb_num = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.2), Inches(1), Inches(0.5))
        tf_num = tb_num.text_frame
        p_num = tf_num.paragraphs[0]
        r_num = p_num.add_run()
        r_num.text = num
        set_font(r_num, size=28, bold=True, color=color)

        # 제목
        tb = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.75), card_w - Inches(0.3), Inches(0.35))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        set_font(r, size=14, bold=True, color=COLOR['text'])

        # 구분선
        sep = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x + Inches(0.15), y + Inches(1.15), card_w - Inches(0.3), Inches(0.02)
        )
        sep.fill.solid()
        sep.fill.fore_color.rgb = COLOR['light_gray']
        sep.line.fill.background()

        # 불릿 아이템
        items_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(1.3), card_w - Inches(0.4), Inches(3))
        itf = items_box.text_frame
        itf.word_wrap = True
        for j, item in enumerate(items):
            ip = itf.paragraphs[0] if j == 0 else itf.add_paragraph()
            ip.space_before = Pt(10)
            ir = ip.add_run()
            ir.text = f"  {item}"
            set_font(ir, size=11, color=COLOR['text'])


def slide_16_data_model(prs):
    """4개 핵심 테이블로 데이터 무결성을 보장합니다"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "4개 핵심 테이블로 데이터 무결성을 보장합니다")
    add_footer(slide, 16)

    tables_data = [
        ("invoices", "인보이스 마스터", COLOR['primary'],
         ["id (PK)", "vendor_id (FK)", "invoice_number", "amount", "status",
          "issue_date", "due_date"]),
        ("invoice_items", "인보이스 항목", COLOR['secondary'],
         ["id (PK)", "invoice_id (FK)", "description", "quantity", "unit_price",
          "amount", "tax_rate"]),
        ("vendors", "거래처 마스터", COLOR['accent'],
         ["id (PK)", "name", "business_number", "contact_email", "payment_terms",
          "bank_account", "status"]),
        ("approvals", "결재 이력", COLOR['success'],
         ["id (PK)", "invoice_id (FK)", "approver_id (FK)", "status", "level",
          "approved_at", "comment"]),
    ]

    card_w = Inches(2.85)
    card_h = Inches(4.8)
    start_x = Inches(0.5)
    gap = Inches(0.2)
    y = CONTENT_TOP + Inches(0.15)

    for i, (table_name, display_name, color, fields) in enumerate(tables_data):
        x = start_x + i * (card_w + gap)

        # 카드 배경
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR['white']
        card.line.color.rgb = color
        card.line.width = Pt(2)
        card.adjustments[0] = 0.04

        # 테이블 헤더 (ER 스타일)
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y, card_w, Inches(0.8)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = color
        header.line.fill.background()

        # 테이블 이름
        tb = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.05), card_w - Inches(0.2), Inches(0.35))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = table_name
        set_font(r, size=14, bold=True, color=COLOR['white'], name=FONT_EN)

        tb2 = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.42), card_w - Inches(0.2), Inches(0.3))
        tf2 = tb2.text_frame
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = display_name
        set_font(r2, size=11, color=RGBColor(0xDD, 0xDD, 0xFF))

        # 필드 목록
        fields_box = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.95), card_w - Inches(0.3), Inches(3.7))
        ftf = fields_box.text_frame
        ftf.word_wrap = True
        for j, field in enumerate(fields):
            fp = ftf.paragraphs[0] if j == 0 else ftf.add_paragraph()
            fp.space_before = Pt(4)
            fr = fp.add_run()
            is_pk = "(PK)" in field
            is_fk = "(FK)" in field
            fr.text = f"  {field}"
            fc = COLOR['danger'] if is_pk else (COLOR['secondary'] if is_fk else COLOR['text'])
            set_font(fr, size=10, bold=(is_pk or is_fk), color=fc, name=FONT_EN)

        # 테이블 간 관계 화살표
        if i < 3:
            arr = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                x + card_w + Inches(0.02), y + Inches(1.5),
                Inches(0.16), Inches(0.3)
            )
            arr.fill.solid()
            arr.fill.fore_color.rgb = COLOR['light_gray']
            arr.line.fill.background()


def slide_17_qa(prs):
    """Q & A 마무리"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 네이비 풀 배경
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR['primary']
    bg.line.fill.background()

    # 상단 골드 라인
    top_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.08)
    )
    top_line.fill.solid()
    top_line.fill.fore_color.rgb = COLOR['accent']
    top_line.line.fill.background()

    # Q & A 텍스트
    txBox = slide.shapes.add_textbox(Inches(2), Inches(2.2), Inches(9), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Q & A"
    set_font(r, size=52, bold=True, color=COLOR['white'])

    # 골드 구분선
    gold_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.8), Inches(2.3), Inches(0.05)
    )
    gold_line.fill.solid()
    gold_line.fill.fore_color.rgb = COLOR['accent']
    gold_line.line.fill.background()

    # 서브 텍스트
    sub = slide.shapes.add_textbox(Inches(2), Inches(4.1), Inches(9), Inches(0.6))
    stf = sub.text_frame
    sp = stf.paragraphs[0]
    sp.alignment = PP_ALIGN.CENTER
    sr = sp.add_run()
    sr.text = "AI 인보이스 자동화 시스템으로 업무 혁신을 시작하세요"
    set_font(sr, size=18, color=COLOR['accent'])

    # 하단 연락처
    contact = slide.shapes.add_textbox(Inches(2), Inches(5.2), Inches(9), Inches(1))
    ctf = contact.text_frame
    lines = ["전략 기획팀  |  AI Innovation Division", "ai-invoice@company.com"]
    for i, line in enumerate(lines):
        cp = ctf.paragraphs[0] if i == 0 else ctf.add_paragraph()
        cp.alignment = PP_ALIGN.CENTER
        cp.space_before = Pt(6)
        cr = cp.add_run()
        cr.text = line
        set_font(cr, size=13, color=COLOR['light_gray'])

    # 하단 골드 라인
    bottom_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08)
    )
    bottom_line.fill.solid()
    bottom_line.fill.fore_color.rgb = COLOR['accent']
    bottom_line.line.fill.background()

    # 장식 원
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(-1), Inches(4.5), Inches(3), Inches(3)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0x25, 0x50, 0x80)
    circle.line.fill.background()

    circle2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(11.5), Inches(0.5), Inches(2.5), Inches(2.5)
    )
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = RGBColor(0x25, 0x50, 0x80)
    circle2.line.fill.background()


# ─── 메인 ────────────────────────────────────────────────────────────

def main():
    print("컨설턴트급 PPT 생성 시작...")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_builders = [
        slide_01_cover,
        slide_02_problem,
        slide_03_solution,
        slide_04_ocr_comparison,
        slide_05_validation,
        slide_06_value_prop,
        slide_07_as_is_to_be,
        slide_08_roi,
        slide_09_approval,
        slide_10_build_vs_buy,
        slide_11_ai_erp,
        slide_12_architecture,
        slide_13_security,
        slide_14_roadmap,
        slide_15_differentiators,
        slide_16_data_model,
        slide_17_qa,
    ]

    for i, builder in enumerate(slide_builders, 1):
        print(f"  슬라이드 {i}/17: {builder.__doc__[:30]}...")
        builder(prs)

    output_path = os.path.join(os.path.dirname(__file__), 'AI_Invoice_Consultant_Grade.pptx')
    prs.save(output_path)
    print(f"\n생성 완료: {output_path}")
    print(f"슬라이드 수: {len(prs.slides)}")

    # 임시 차트 파일 정리
    for f in temp_files:
        try:
            os.remove(f)
        except OSError:
            pass
    print("임시 파일 정리 완료")


if __name__ == '__main__':
    main()
